"""
Generate a 3-slide weekly update deck:
  - Slides 1–2: last week progress (with quantitative highlights)
  - Slide 3: next steps
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Weekly_Update_Last_Week.pptx")

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCA, 0xCA, 0xDA)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.35), left=Inches(0.55)):
    box = slide.shapes.add_textbox(left, top, Inches(8.9), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def add_subtitle(slide, text, top=Inches(1.05), left=Inches(0.55)):
    box = slide.shapes.add_textbox(left, top, Inches(8.9), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT


def add_bullets(
    slide,
    items,
    top=Inches(1.55),
    left=Inches(0.6),
    width=Inches(8.8),
    font_pt=15,
):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_pt)
        run.font.color.rgb = color if color else WHITE


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(9), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, top, left=Inches(0.55)):
    nrows, ncols = len(rows) + 1, len(headers)
    h = min(Inches(0.32 * nrows + 0.15), Inches(1.8))
    shape = slide.shapes.add_table(nrows, ncols, left, top, Inches(8.9), h)
    tbl = shape.table
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = htxt
        for para in c.text_frame.paragraphs:
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.text = str(val)
            for para in c.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0x24, 0x24, 0x3E)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    footer = (
        "Efficient Transformer Inference on Commodity GPUs  |  "
        "Jay Darji & Karma Patel  |  SJSU Spring 2026"
    )

    # --- Slide 1: Progress (wrap-up / documentation) ---
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, DARK_BG)
    add_title(s1, "Weekly Update — Last Week Progress")
    add_subtitle(s1, "Final submission alignment, repository, and documentation")
    add_bullets(
        s1,
        [
            ("Completed Week 14 deliverables: final report, explanation PDF, reproducible codebase.", WHITE),
            ("", None),
            ("Repository published (code only; results excluded via .gitignore for size).", WHITE),
            ("README expanded: archived results layout, 14 required plots (150 DPI), key JSON data paths.", WHITE),
            ("", None),
            ("Hardware & model held constant for comparability:", YELLOW),
            ("    NVIDIA RTX 3050 Ti Laptop (4 GB VRAM)  ·  Qwen2.5-3B-Instruct (4-bit NF4)", LIGHT),
        ],
        top=Inches(1.5),
    )
    add_footer(s1, footer)

    # --- Slide 2: Quantitative progress (measured outcomes) ---
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, DARK_BG)
    add_title(s2, "Last Week — Quantitative Highlights")
    add_subtitle(s2, "OOM threshold & quality (greedy decode, batch size 1, fixed benchmark protocol)")
    add_bullets(
        s2,
        [
            ("Max context before OOM (long-context sweep):", YELLOW),
            ("    Baseline (eager): 6,784 tokens", WHITE),
            ("    SDPA default: 7,040 tokens  (+3.8% vs baseline)", GREEN),
            ("    FlashAttention-2: 8,064 tokens  (+18.9% vs baseline)", GREEN),
            ("", None),
            ("Quality retention: ≤2% accuracy drop vs baseline on ARC-Easy, BoolQ, HellaSwag (200 ex. each).", WHITE),
            ("", None),
            ("KV INT4: lower KV memory per token; combined attention + KV quant unstable on this GQA setup (documented).", LIGHT),
        ],
        top=Inches(1.45),
    )
    add_table(
        s2,
        ["Configuration", "OOM threshold", "Δ vs baseline"],
        [
            ["Eager baseline", "6,784 tok", "—"],
            ["SDPA default", "7,040 tok", "+3.8%"],
            ["FlashAttention-2", "8,064 tok", "+18.9%"],
        ],
        top=Inches(4.85),
    )
    add_footer(s2, footer)

    # --- Slide 3: Next phase — combined KV + FlashAttention on second model ---
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, DARK_BG)
    add_title(s3, "Next Phase — Combined Performance")
    add_subtitle(
        s3,
        "Validate FlashAttention-2 + quantized KV-cache together on a multi-head attention (MHA) model",
    )
    add_bullets(
        s3,
        [
            ("Why a new model: Qwen2.5-3B uses GQA; combined attention + QuantizedCache was unstable or kernel-limited.", LIGHT),
            ("", None),
            ("Target model (HF):", YELLOW),
            ("    microsoft/phi-2  (Phi-2, ~2.7B) — standard MHA, fits existing 4-bit NF4 + transformers stack.", GREEN),
            ("", None),
            ("Experiment plan:", YELLOW),
            ("    Same protocol as weeks 9–10: attn_implementation=flash_attention_2 (or SDPA) + QuantizedCache INT4 → INT2.", WHITE),
            ("    Metrics: OOM threshold, p50/p95 latency, tokens/s, peak VRAM, quality vs FP16 KV baseline.", WHITE),
            ("", None),
            ("Optional fast debug: TinyLlama/TinyLlama-1.1B-Chat-v1.0 before full Phi-2 runs.", LIGHT),
        ],
        top=Inches(1.42),
        font_pt=13,
    )
    add_footer(s3, footer)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
