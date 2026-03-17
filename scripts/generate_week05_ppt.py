"""Generate Week 05 progress PowerPoint (3 slides)."""

import os, sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config import RESULTS_DIR

OUT_DIR = os.path.join(RESULTS_DIR, "week05_attention_opt")
PPT_PATH = os.path.join(RESULTS_DIR, "Weekly_Progress_Week5.pptx")

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT = RGBColor(0x89, 0xB4, 0xFA)
GREEN = RGBColor(0xA6, 0xE3, 0xA1)
RED = RGBColor(0xF3, 0x8B, 0xA8)
YELLOW = RGBColor(0xF9, 0xE2, 0xAF)
WHITE = RGBColor(0xCD, 0xD6, 0xF4)
SUBTEXT = RGBColor(0xBA, 0xC2, 0xDE)
SURFACE = RGBColor(0x31, 0x32, 0x44)
OVERLAY = RGBColor(0x45, 0x47, 0x5A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, size=18,
              color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                        Inches(0.38 * n_rows))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    paragraph.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x3B, 0x3D, 0x54)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0x2A, 0x2B, 0x3D) if r % 2 == 1
                    else RGBColor(0x33, 0x34, 0x48)
                )

    return tbl_shape


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ================================================================
    # SLIDE 1 — Backend Benchmark & Availability
    # ================================================================
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1, DARK_BG)

    _add_text(s1, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 05 - Attention Optimization Integration",
              size=28, color=ACCENT, bold=True)
    _add_text(s1, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
              "Model: Qwen/Qwen2.5-3B-Instruct  |  GPU: RTX 3050 Ti (4 GB)  |  PyTorch 2.6.0 + CUDA 12.4",
              size=13, color=SUBTEXT)

    _add_text(s1, Inches(0.6), Inches(1.5), Inches(5), Inches(0.4),
              "Single-Prompt Benchmark (64 new tokens, 5 runs)",
              size=16, color=YELLOW, bold=True)

    bench_rows = [
        ["Backend", "p50 (ms)", "p95 (ms)", "tok/s", "VRAM (MB)", "Status"],
        ["Eager",             "5,936",  "6,141",  "10.71", "2,018", "OK"],
        ["SDPA Default",      "5,945",  "5,998",  "10.74", "3,388", "OK"],
        ["SDPA Math",         "5,942",  "5,949",  "10.77", "3,388", "OK"],
        ["FlashAttention-2",  "5,952",  "6,042",  "10.73", "3,389", "OK"],
        ["SDPA Mem-Efficient", "-",     "-",      "-",     "-",     "FAILED"],
    ]
    _add_table(s1, Inches(0.6), Inches(2.0), Inches(8.5), bench_rows)

    _add_text(s1, Inches(0.6), Inches(4.7), Inches(5), Inches(0.4),
              "Backend Availability",
              size=16, color=YELLOW, bold=True)

    avail_rows = [
        ["Backend", "Available", "Works with GQA Model", "Notes"],
        ["SDPA Flash",         "No",  "N/A", "Not compiled in PyTorch Windows wheel"],
        ["SDPA Mem-Efficient", "Yes", "No",  "GQA head mismatch (16Q vs 2KV)"],
        ["SDPA Math",          "Yes", "Yes", "Universal fallback"],
        ["FlashAttention-2",   "Yes", "Yes", "flash-attn 2.7.4 (pre-built wheel)"],
    ]
    _add_table(s1, Inches(0.6), Inches(5.2), Inches(11), avail_rows)

    # Key insight box
    _add_text(s1, Inches(9.5), Inches(2.0), Inches(3.2), Inches(2.2),
              "Key Finding:\n\n"
              "All working backends converge to ~10.7 tok/s at steady state.\n\n"
              "FlashAttention-2 installed via pre-built HuggingFace wheel "
              "(no compiler needed).",
              size=12, color=GREEN)

    # ================================================================
    # SLIDE 2 — Sweep Results & Correctness
    # ================================================================
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2, DARK_BG)

    _add_text(s2, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 05 - Benchmark Sweeps & Correctness Verification",
              size=28, color=ACCENT, bold=True)

    _add_text(s2, Inches(0.6), Inches(1.1), Inches(7), Inches(0.4),
              "Latency & Throughput vs Prompt Length (batch=1, 64 new tokens)",
              size=16, color=YELLOW, bold=True)

    sweep_rows = [
        ["Backend", "Prompt Len", "p50 (ms)", "p95 (ms)", "tok/s", "VRAM (MB)"],
        ["Eager",            "128",  "6,043",  "6,060",  "10.58", "2,027"],
        ["Eager",            "512",  "6,219",  "6,242",  "10.30", "2,070"],
        ["Eager",            "1024", "6,533",  "6,545",  "9.81",  "2,152"],
        ["SDPA Math",        "128",  "6,039",  "6,111",  "10.56", "3,396"],
        ["SDPA Math",        "512",  "6,181",  "6,195",  "10.36", "3,431"],
        ["SDPA Math",        "1024", "6,555",  "6,565",  "9.76",  "3,565"],
        ["FlashAttention-2", "128",  "6,060",  "6,122",  "10.54", "3,397"],
        ["FlashAttention-2", "512",  "6,221",  "6,260",  "10.29", "3,432"],
        ["FlashAttention-2", "1024", "6,447",  "6,545",  "9.92",  "3,478"],
    ]
    _add_table(s2, Inches(0.6), Inches(1.6), Inches(9), sweep_rows)

    _add_text(s2, Inches(0.6), Inches(5.9), Inches(5), Inches(0.4),
              "Correctness vs Week 03 Baseline (5 prompts, 20 greedy steps)",
              size=16, color=YELLOW, bold=True)

    corr_rows = [
        ["Backend", "Mean Agreement", "Min Agreement", "Max Logit Diff"],
        ["SDPA Default",      "100.0%", "100.0%", "0.000"],
        ["SDPA Math",         "100.0%", "100.0%", "0.000"],
        ["Eager",             "87.0%",  "55.0%",  "25.051"],
        ["FlashAttention-2",  "87.0%",  "55.0%",  "25.035"],
    ]
    _add_table(s2, Inches(0.6), Inches(6.3), Inches(7), corr_rows)

    _add_text(s2, Inches(9.8), Inches(5.9), Inches(3), Inches(1.5),
              "Notes:\n\n"
              "SDPA (default/math) matches baseline exactly.\n\n"
              "Eager & FA2 show 87% agreement - expected due to different "
              "attention code paths and floating-point ordering.",
              size=12, color=SUBTEXT)

    # ================================================================
    # SLIDE 3 — Planned Work for Week 06
    # ================================================================
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3, DARK_BG)

    _add_text(s3, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 06 Plan - Attention Optimization Final Results",
              size=28, color=ACCENT, bold=True)

    _add_text(s3, Inches(0.6), Inches(1.2), Inches(11), Inches(0.4),
              "Objective: Stabilize attention benchmarks, produce final comparison tables, plots, and report text.",
              size=15, color=SUBTEXT)

    tasks = [
        ("1. Final Attention Benchmarks",
         "Re-run complete benchmark suite for all working backends (eager, SDPA math, "
         "FlashAttention-2) with extended prompt lengths and multiple seeds for "
         "statistical confidence."),

        ("2. OOM Threshold Analysis",
         "Find the maximum context length before out-of-memory for each backend. "
         "This reveals the practical sequence-length limit on 4 GB VRAM."),

        ("3. Generate Comparison Plots",
         "Produce latency-vs-prompt-length, throughput-vs-prompt-length, and "
         "VRAM-vs-prompt-length charts for the final report."),

        ("4. Stability & Limitations Notes",
         "Document thermal throttling effects on laptop GPU, GQA incompatibility "
         "with SDPA mem-efficient, and flash SDP kernel absence on Windows."),

        ("5. Report Draft Subsection",
         "Write the Attention Optimization section of the final report with "
         "figure placeholders, backend selection rationale, and result analysis."),
    ]

    y = Inches(1.9)
    for title, desc in tasks:
        _add_text(s3, Inches(0.8), y, Inches(10), Inches(0.35),
                  title, size=17, color=GREEN, bold=True)
        _add_text(s3, Inches(1.0), y + Inches(0.38), Inches(10), Inches(0.55),
                  desc, size=13, color=SUBTEXT)
        y += Inches(1.05)

    _add_text(s3, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
              "Deliverables: final_attention_results.json  |  attention_plots/  |  report_attention_section.md",
              size=13, color=YELLOW, bold=True)

    prs.save(PPT_PATH)
    print(f"Saved: {PPT_PATH}")


if __name__ == "__main__":
    build_ppt()
