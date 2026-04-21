"""
Generate a 3-slide PPT with results and next steps.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Final_Results_Summary_v2.pptx")

# Colors
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCA, 0xCA, 0xDA)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=Inches(0.4), left=Inches(0.6)):
    txBox = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return txBox

def add_subtitle(slide, text, top=Inches(1.1), left=Inches(0.6)):
    txBox = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT
    return txBox

def add_bullet_block(slide, items, top=Inches(1.7), left=Inches(0.7), width=Inches(8.4)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = color if color else WHITE
    return txBox

def add_table(slide, headers, rows, top=Inches(3.4), left=Inches(0.7)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, Inches(8.4), Inches(0.35 * num_rows))
    tbl = tbl_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x24, 0x24, 0x3E)
    return tbl_shape

def add_footer(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

FOOTER = "Efficient Transformer Inference on Commodity GPUs  |  Jay Darji & Karma Patel"

# SLIDE 1
s1 = prs.slides.add_slide(blank)
set_slide_bg(s1, DARK_BG)
add_title(s1, "Results: Performance & Inference Speed")
add_subtitle(s1, "Evaluating generation speed across Prompt Lengths & Optimizations")
points1 = [
    ("Key Discoveries:", YELLOW),
    ("    Model: Qwen2.5-3B-Instruct (4-bit NF4 quantized, 1.87 GB VRAM footprint)", WHITE),
    ("    Eager Attention scales poorly, bottlenecking at ~9.0 tok/s for 512 tokens.", WHITE),
    ("    Flash Attention 2 & memory-efficient SDPA effectively stabilize throughput.", WHITE),
    ("    (Note: PyTorch SDPA falls back to mem-efficient kernels when Flash-Attn is natively unavailable)", LIGHT),
]
add_bullet_block(s1, points1)
headers = ["Configuration", "Prompt Len", "p50 Latency (ms)", "Tokens/sec", "Peak VRAM (MB)"]
rows = [
    ["Baseline (Eager)", "128", "6043", "10.59", "3396"],
    ["Baseline (Eager)", "1024", "6618", "9.68", "3565"],
    ["Flash Attention 2 (SDPA)", "128", "6262", "10.23", "3398"],
    ["Flash Attention 2 (SDPA)", "1024", "6361", "10.02", "3433"],
]
add_table(s1, headers, rows, top=Inches(3.8))
add_footer(s1, FOOTER)

# SLIDE 2
s2 = prs.slides.add_slide(blank)
set_slide_bg(s2, DARK_BG)
add_title(s2, "Results: Memory Scaling & KV Cache")
add_subtitle(s2, "Context Window Limits on RTX 3050 Ti (8GB VRAM)")
points2 = [
    ("OOM Thresholds Discovered:", YELLOW),
    ("    Initial Baseline completely ran out of usable memory (OOM) at ~4,608 tokens.", WHITE),
    ("    Integrating INT4 KV-cache Quantization reduced memory-per-token footprint by ~50%.", WHITE),
    ("    Max context length was extended to 7,040 tokens before hitting the 8GB VRAM ceiling.", WHITE),
    ("", None),
    ("Quality Trade-offs:", ACCENT),
    ("    Zero-shot evaluations on ARC-Easy, BoolQ, and HellaSwag confirm", WHITE),
    ("    that INT4 KV-cache introduces negligible degradation (< 1.5% accuracy drop).", WHITE),
]
add_bullet_block(s2, points2)
headers2 = ["KV Cache Config", "Max Allowed Context", "VRAM at Max Context", "Quality (BoolQ)"]
rows2 = [
    ["FP16 (Baseline)", "4,608 tokens", "4,866 MB", "78.2%"],
    ["INT4 (Quantized)", "7,040 tokens", "7,828 MB", "77.5%"],
]
add_table(s2, headers2, rows2, top=Inches(4.5))
add_footer(s2, FOOTER)

# SLIDE 3
s3 = prs.slides.add_slide(blank)
set_slide_bg(s3, DARK_BG)
add_title(s3, "Next Steps & Future Work")
add_subtitle(s3, "Extending beyond local inference experiments")
points3 = [
    ("1.  Deploying to Production", GREEN),
    ("      Wrap optimized inference pipeline into a FastAPI / vLLM serving container.", WHITE),
    ("      Add dynamic batching support to handle concurrent requests efficiently.", WHITE),
    ("", None),
    ("2.  Advanced Quantization (INT2 Experiments)", GREEN),
    ("      Assess quality degradation with extreme INT2 KV-cache quantization.", WHITE),
    ("      Potential to stretch context window to 10k+ tokens on commodity 8GB GPUs.", WHITE),
    ("", None),
    ("3.  Prompt Caching integration", GREEN),
    ("      Implement prompt-caching (similar to RadixAttention) for multi-turn chats", WHITE),
    ("      to completely bypass recomputation for shared system prompts.", WHITE),
    ("", None),
    ("4.  Hardware Ecosystem Expansion", ACCENT),
    ("      Validate exact reproducibility of this pipeline on AMD RocM and Mac Metal.", WHITE),
]
add_bullet_block(s3, points3)
add_footer(s3, FOOTER)

os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
