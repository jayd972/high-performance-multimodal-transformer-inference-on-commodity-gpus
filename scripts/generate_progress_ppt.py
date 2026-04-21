"""
Generate a 3-slide progress deck:
  Slide 1 — Combined experiments: what we did & stability results
  Slide 2 — Quantitative benchmarks with tables (Phi-2 & TinyLlama)
  Slide 3 — Future work: multimodal model (LLaVA-Phi-2)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Progress_Update_Combined_Experiments.pptx")

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCA, 0xCA, 0xDA)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
YELLOW  = RGBColor(0xF1, 0xC4, 0x0F)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
CELL_BG = RGBColor(0x24, 0x24, 0x3E)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.3), left=Inches(0.5)):
    box = slide.shapes.add_textbox(left, top, Inches(9.0), Inches(0.7))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = ACCENT


def add_subtitle(slide, text, top=Inches(0.95), left=Inches(0.5)):
    box = slide.shapes.add_textbox(left, top, Inches(9.0), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(12); p.font.color.rgb = LIGHT


def add_bullets(slide, items, top=Inches(1.45), left=Inches(0.55), width=Inches(8.9),
                font_pt=14):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.3))
    tf = box.text_frame; tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_pt)
        run.font.color.rgb = color if color else WHITE


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(9), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(9); p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER


def add_table(slide, headers, rows, top, left=Inches(0.5), width=Inches(9.0),
              row_height=Inches(0.28), header_font=10, cell_font=9):
    nrows, ncols = len(rows) + 1, len(headers)
    h = row_height * (nrows + 0.4)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, h)
    tbl = shape.table
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = htxt
        for para in c.text_frame.paragraphs:
            para.font.size = Pt(header_font); para.font.bold = True
            para.font.color.rgb = WHITE; para.alignment = PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j); c.text = str(val)
            for para in c.text_frame.paragraphs:
                para.font.size = Pt(cell_font); para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER
            c.fill.solid(); c.fill.fore_color.rgb = CELL_BG


def add_label(slide, text, top, left, width=Inches(9.0), size=11, color=YELLOW,
              bold=True):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    footer = (
        "Efficient Transformer Inference on Commodity GPUs  |  "
        "Jay Darji & Karma Patel  |  SJSU Spring 2026"
    )

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1 — What we did & stability
    # ═══════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, DARK_BG)
    add_title(s1, "Combined Attention + KV-Cache Experiments")
    add_subtitle(
        s1,
        "FlashAttention-2 / SDPA + QuantizedCache (INT4, INT2) on MHA models  —  RTX 3050 Ti (4 GB)",
    )
    add_bullets(
        s1,
        [
            ("Problem: Qwen2.5-3B uses GQA → combined attention + KV-cache quantization was unstable.", WHITE),
            ("", None),
            ("Solution: switch to MHA models where kernels compose cleanly:", YELLOW),
            ("    • microsoft/phi-2  (2.7B, MHA 32Q/32KV, 4-bit = 1.66 GB)", GREEN),
            ("    • TinyLlama-1.1B-Chat-v1.0  (1.1B, MHA 32Q/32KV, 4-bit = 0.70 GB)", GREEN),
            ("", None),
            ("Benchmark matrix — 8 configs × 2 models × 3 prompt lengths (128 / 512 / 1024 tokens):", YELLOW),
            ("    Baseline (eager), SDPA, FlashAttention-2, KV INT4,", WHITE),
            ("    Combined SDPA+INT4, Combined SDPA+INT2, Combined Flash+INT4, Combined Flash+INT2", WHITE),
            ("", None),
            ("Result: all 4 combined configs stable on BOTH models (4/4 each).", GREEN),
            ("    OOM threshold: 8064 tokens across all configs on both models.", WHITE),
        ],
        top=Inches(1.4), font_pt=13,
    )

    add_table(
        s1,
        ["Combined Config", "Phi-2 Stable?", "TinyLlama Stable?"],
        [
            ["SDPA + KV INT4", "✓  YES", "✓  YES"],
            ["SDPA + KV INT2", "✓  YES", "✓  YES"],
            ["Flash + KV INT4", "✓  YES", "✓  YES"],
            ["Flash + KV INT2", "✓  YES", "✓  YES"],
        ],
        top=Inches(5.7), left=Inches(2.0), width=Inches(6.0),
    )
    add_footer(s1, footer)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2 — Quantitative benchmarks
    # ═══════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, DARK_BG)
    add_title(s2, "Benchmark Results — Latency, VRAM & Quality")
    add_subtitle(
        s2,
        "4-bit NF4, batch=1, 1024-token prompt, 64-token output, greedy decode  —  5 runs each",
    )

    add_label(s2, "Phi-2 (2.7B)  —  Peak VRAM & Latency at 1024 tokens",
              top=Inches(1.35), left=Inches(0.5))
    add_table(
        s2,
        ["Config", "p50 (ms)", "Tok/s", "VRAM (MB)", "ΔVRAM"],
        [
            ["Baseline (eager)", "5,432", "11.7", "2,361", "—"],
            ["SDPA default", "4,719", "13.4", "3,409", "+44%"],
            ["KV INT4 (eager)", "6,300", "10.2", "2,156", "−9%"],
            ["SDPA + KV INT4", "6,025", "10.6", "1,963", "−17%"],
            ["Flash + KV INT2", "1,107", "57.7", "1,899", "−20%"],
        ],
        top=Inches(1.65), left=Inches(0.5), width=Inches(9.0),
        row_height=Inches(0.26), header_font=9, cell_font=9,
    )

    add_label(s2, "TinyLlama (1.1B)  —  Peak VRAM & Latency at 1024 tokens",
              top=Inches(3.7), left=Inches(0.5))
    add_table(
        s2,
        ["Config", "p50 (ms)", "Tok/s", "VRAM (MB)", "ΔVRAM"],
        [
            ["Baseline (eager)", "4,165", "15.4", "1,051", "—"],
            ["SDPA default", "5,099", "12.5", "1,598", "+52%"],
            ["KV INT4 (eager)", "5,651", "11.3", "1,034", "−2%"],
            ["Flash + KV INT4", "5,433", "11.8", "814", "−23%"],
            ["Flash + KV INT2", "5,515", "11.6", "811", "−23%"],
        ],
        top=Inches(4.0), left=Inches(0.5), width=Inches(9.0),
        row_height=Inches(0.26), header_font=9, cell_font=9,
    )

    add_label(s2, "Quality Retention (0-shot, 200 examples/dataset)",
              top=Inches(5.9), left=Inches(0.5))
    add_table(
        s2,
        ["Model", "Config", "ARC-Easy", "BoolQ", "HellaSwag", "Mean"],
        [
            ["Phi-2", "Comb. SDPA+INT4", "77.5%", "60.5%", "68.5%", "68.8%"],
            ["Phi-2", "Comb. Flash+INT2", "77.5%", "60.5%", "68.5%", "68.8%"],
            ["TinyLlama", "Baseline", "43.5%", "41.5%", "54.5%", "46.5%"],
            ["TinyLlama", "Comb. Flash+INT2", "43.5%", "42.0%", "54.5%", "46.7%"],
        ],
        top=Inches(6.15), left=Inches(0.5), width=Inches(9.0),
        row_height=Inches(0.24), header_font=9, cell_font=8,
    )
    add_footer(s2, footer)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3 — Future work: multimodal
    # ═══════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, DARK_BG)
    add_title(s3, "Future Work — Multimodal Inference")
    add_subtitle(
        s3,
        "Extend combined optimizations to a vision-language model on the same 4 GB GPU",
    )
    add_bullets(
        s3,
        [
            ("Target model:", YELLOW),
            ("    xtuner/llava-phi-2-3b-hf  (LLaVA-Phi-2, ~3B)", GREEN),
            ("    Phi-2 language backbone + SigLIP vision encoder, MIT license", WHITE),
            ("", None),
            ("Why this model:", YELLOW),
            ("    MHA (32Q/32KV) — we already proved Phi-2 composes with Flash+KV quant.", WHITE),
            ("    ~1.7 GB in 4-bit NF4 — fits within 4 GB VRAM budget with headroom for vision.", WHITE),
            ("    Multimodal: image understanding + text generation (not just text-to-text).", WHITE),
            ("    HuggingFace transformers-native (LlavaForConditionalGeneration).", WHITE),
            ("", None),
            ("Experiment plan:", YELLOW),
            ("    1. Baseline: LLaVA-Phi-2 4-bit, eager attention, FP16 KV.", WHITE),
            ("    2. Combined: flash_attention_2 + QuantizedCache INT4/INT2.", WHITE),
            ("    3. Metrics: latency, tokens/s, peak VRAM (text + vision), OOM threshold.", WHITE),
            ("    4. Quality: VQA accuracy on a small eval set vs baseline.", WHITE),
            ("", None),
            ("Smaller alternative: bczhou/tiny-llava-v1-hf (~1.4B, TinyLlama + CLIP vision).", LIGHT),
        ],
        top=Inches(1.35), font_pt=12,
    )
    add_footer(s3, footer)

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
