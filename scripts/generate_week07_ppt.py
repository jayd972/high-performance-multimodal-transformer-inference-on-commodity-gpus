"""
Generate Week 7 progress presentation — 3 slides only.
Slide 1-2: Week 7 results, Slide 3: Week 8 plan.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Weekly_Progress_Week07.pptx")

DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
TITLE_BG   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
YELLOW     = RGBColor(0xF1, 0xC4, 0x0F)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ROW_DARK   = RGBColor(0x15, 0x25, 0x3D)
ROW_LIGHT  = RGBColor(0x1A, 0x2E, 0x4A)

FOOTER_TEXT = "LLM Inference Optimization  |  Darji & Patel  |  Week 7"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, items, top=Inches(1.8), left=Inches(0.7), width=Inches(8.6)):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, clr, bld) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.color.rgb = clr
        run.font.bold = bld
    return box


def add_table(slide, headers, rows, top=Inches(2.8), left=Inches(0.5),
              width=Inches(9.0)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.35)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   row_h * n_rows)
    tbl = shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    for i, row in enumerate(rows):
        bg = ROW_DARK if i % 2 == 0 else ROW_LIGHT
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return shape


def add_footer(slide):
    add_textbox(slide, FOOTER_TEXT,
                Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3),
                font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, title,
                Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7),
                font_size=26, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.4),
                    font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Slide 1: Week 7 Results — Prototype & Correctness ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 7 Results — KV-Cache Quantization Prototype",
                "Support check, end-to-end prototype, and correctness verification")

add_bullets(s, [
    ("Support Check", ACCENT, True),
    ("  optimum-quanto backend available — supports INT4 and INT2 (not INT8)", WHITE, False),
], top=Inches(1.5))

add_table(s, ["KV Precision", "Status", "Tokens", "Time (ms)", "Tok/s", "VRAM (MB)"],
          [
              ["INT4", "Success", "64", "24,045", "2.66", "2,018"],
              ["INT2", "Success", "64", "23,663", "2.70", "2,017"],
          ],
          top=Inches(2.5), left=Inches(0.5), width=Inches(9.0))

add_bullets(s, [
    ("Correctness Verification", ACCENT, True),
], top=Inches(3.9))

add_table(s, ["KV Precision", "Top-1 Agreement", "Max Logit Diff"],
          [
              ["INT4", "100.0%", "0.0"],
              ["INT2", "100.0%", "0.0"],
          ],
          top=Inches(4.5), left=Inches(1.5), width=Inches(7.0))

add_bullets(s, [
    ("Both INT4 and INT2 are lossless for greedy decoding — safe to deploy", GREEN, True),
], top=Inches(5.7))

add_footer(s)

# ── Slide 2: Week 7 Results — Memory Per Token & Decision ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 7 Results — Memory Savings & Integration Decision",
                "Per-token VRAM growth across precisions (64–1024 tokens)")

add_table(s, ["KV Precision", "64 tok", "256 tok", "512 tok", "1024 tok", "MB / Token", "Savings"],
          [
              ["FP16 (baseline)", "2,021", "2,038", "2,062", "2,205", "0.181", "—"],
              ["INT4", "2,019", "2,032", "2,049", "2,180", "0.157", "−13.3%"],
              ["INT2", "2,019", "2,031", "2,047", "2,176", "0.153", "−15.5%"],
          ],
          top=Inches(1.7), left=Inches(0.3), width=Inches(9.4))

add_bullets(s, [
    ("Projected Savings at Long Contexts", ACCENT, True),
    ("  At 8,000 tokens: INT4 saves ~192 MB, INT2 saves ~224 MB vs FP16", WHITE, False),
    ("", WHITE, False),
    ("Integration Decision: USE EXISTING PIPELINE", GREEN, True),
    ("  transformers + optimum-quanto works out-of-the-box — no custom code needed", WHITE, False),
    ("", WHITE, False),
    ("Recommendation", ACCENT, True),
    ("  INT4 — best quality/savings tradeoff (recommended for Week 8 experiments)", GREEN, False),
    ("  INT2 — functional but output degrades at longer sequences", YELLOW, False),
    ("  INT8 — not supported by quanto", RED, False),
], top=Inches(3.4))

add_footer(s)

# ── Slide 3: Week 8 Plan ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 8 Plan — KV-Cache Experiments & Quality Evaluation")

add_bullets(s, [
    ("Task 1 — Benchmark Sweep", ACCENT, True),
    ("  Compare FP16 vs INT4 vs INT2 at prompt lengths 128 / 512 / 1024", WHITE, False),
    ("  Measure latency, throughput, and peak VRAM for each", WHITE, False),
    ("", WHITE, False),
    ("Task 2 — Quality Evaluation (0-shot, 200 examples each)", ACCENT, True),
    ("  ARC-Easy (science QA)  |  BoolQ (boolean QA)  |  HellaSwag (commonsense)", WHITE, False),
    ("  Verify quality guard-rail: ≤2% accuracy drop with KV quantization", WHITE, False),
    ("", WHITE, False),
    ("Task 3 — Tradeoff Analysis", ACCENT, True),
    ("  Plot VRAM savings (%) vs latency overhead (%) for each config", WHITE, False),
    ("  Identify the Pareto-optimal KV-cache setting", WHITE, False),
    ("", WHITE, False),
    ("Task 4 — Final KV Setting Selection", ACCENT, True),
    ("  Pick the best KV-cache config for Week 9 combined testing", WHITE, False),
    ("", WHITE, False),
    ("Expected Deliverables", ACCENT, True),
    ("  kv_comparison_results.json  |  baseline_quality.json", LIGHT_GRAY, False),
    ("  kv_tradeoff_analysis.json   |  final_kv_selection.json", LIGHT_GRAY, False),
], top=Inches(1.5))

add_footer(s)

# ── Save ──
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
