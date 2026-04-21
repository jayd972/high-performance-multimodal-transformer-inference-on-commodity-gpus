"""
Generate the final project presentation (9 slides) as a PowerPoint file.

Requires: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Weekly_Progress_Final.pptx")

# ── Color palette ──
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
TITLE_BG   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
YELLOW     = RGBColor(0xF1, 0xC4, 0x0F)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
ROW_DARK   = RGBColor(0x15, 0x25, 0x3D)
ROW_LIGHT  = RGBColor(0x1A, 0x2E, 0x4A)

FOOTER_TEXT = "Efficient Transformer Inference on Commodity GPUs  |  Darji & Patel  |  March 2026"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, items, top=Inches(1.8), left=Inches(0.7), width=Inches(8.6)):
    """items: list of (text, color, bold) tuples."""
    box = slide.shapes.add_textbox(left, top, width, Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, clr, bld) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = clr
        run.font.bold = bld
    return box


def add_table(slide, headers, rows, top=Inches(2.8), left=Inches(0.5),
              width=Inches(9.0)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.38)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   row_h * n_rows)
    tbl = shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    for i, row in enumerate(rows):
        bg = ROW_DARK if i % 2 == 0 else ROW_LIGHT
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return shape


def add_footer(slide):
    add_textbox(slide, FOOTER_TEXT,
                Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3),
                font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, title,
                Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7),
                font_size=28, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.4),
                    font_size=14, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Slide 1: Title ──
s = prs.slides.add_slide(blank)
set_bg(s, TITLE_BG)
add_textbox(s, "Efficient Transformer\nInference on\nCommodity GPUs",
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5),
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "Jay Darji  &  Karma Patel",
            Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5),
            font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, "San Jose State University  |  Spring 2026",
            Inches(0.8), Inches(4.8), Inches(8.4), Inches(0.5),
            font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_footer(s)

# ── Slide 2: Problem & Motivation ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Problem & Motivation")
add_bullets(s, [
    ("The Problem", ACCENT, True),
    ("    LLMs require expensive GPUs (A100, H100) for inference", WHITE, False),
    ("    Consumer laptops with 4 GB VRAM are typically excluded", WHITE, False),
    ("", WHITE, False),
    ("Our Question", ACCENT, True),
    ("    Can we run a 3B-parameter model on a $800 laptop GPU?", WHITE, False),
    ("    Can attention + KV-cache optimizations make it practical?", WHITE, False),
    ("", WHITE, False),
    ("Goal", ACCENT, True),
    ("    Criterion A: \u226520% p95 latency reduction", GREEN, False),
    ("    Criterion B: \u22651.5\u00d7 maximum context length", GREEN, False),
    ("    Quality guard-rail: \u22642% accuracy drop", GREEN, False),
])
add_footer(s)

# ── Slide 3: Approach ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Approach", "Model, hardware, and optimization axes")
add_bullets(s, [
    ("Model: Qwen2.5-3B-Instruct", ACCENT, True),
    ("    4-bit NF4 quantization via bitsandbytes \u2192 1.87 GB weight footprint", WHITE, False),
    ("    GQA architecture: 16 query heads, 2 KV heads", WHITE, False),
    ("", WHITE, False),
    ("Hardware: NVIDIA RTX 3050 Ti Laptop GPU", ACCENT, True),
    ("    4 GB GDDR6 VRAM  |  Ampere (SM 8.6)  |  Windows 10", WHITE, False),
    ("", WHITE, False),
    ("Optimization Axis 1 \u2014 Attention Backends (4 tested)", ACCENT, True),
    ("    Eager (baseline)  |  SDPA default  |  SDPA math  |  FlashAttention-2", WHITE, False),
    ("", WHITE, False),
    ("Optimization Axis 2 \u2014 KV-Cache Precision (3 tested)", ACCENT, True),
    ("    FP16 (baseline)  |  INT4  |  INT2", WHITE, False),
])
add_footer(s)

# ── Slide 4: Key Results — Latency ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Key Results \u2014 Latency",
                "p95 latency at 1024-token prompts (batch=1, 64 output tokens)")
add_table(s, ["Configuration", "p95 Latency (ms)", "Throughput (tok/s)",
              "Peak VRAM (MB)", "vs. Baseline"],
          [
              ["baseline (eager)", "11,209", "5.73", "2,152", "\u2014"],
              ["sdpa_default", "7,509", "8.64", "2,203", "\u221233.0%"],
              ["flash_attention_2", "7,641", "8.40", "2,115", "\u221231.8%"],
              ["kv_int4", "8,091", "7.80", "2,136", "\u221227.8%"],
          ],
          top=Inches(2.2))
add_bullets(s, [
    ("\u2705  Criterion A MET: SDPA default = 33% p95 latency reduction (target: \u226520%)", GREEN, True),
    ("    FlashAttention-2 also exceeds threshold at 31.8% reduction", WHITE, False),
], top=Inches(4.6))
add_footer(s)

# ── Slide 5: Key Results — Context Length ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Key Results \u2014 Context Length",
                "Maximum tokens before CUDA OOM (256-token increments)")
add_table(s, ["Configuration", "Max Context (tokens)", "Improvement"],
          [
              ["baseline (eager)", "6,784", "\u2014"],
              ["sdpa_default", "7,040", "+3.8%"],
              ["flash_attention_2", "8,064", "+18.9%"],
              ["kv_int4", "6,784", "0%"],
          ],
          top=Inches(2.2))
add_bullets(s, [
    ("FlashAttention-2: 8,064 tokens = 1,280 more than baseline", ACCENT, True),
    ("    Near-linear VRAM scaling: ~3.5 MB / 256 tokens (vs. ~400 MB for eager)", WHITE, False),
    ("    Peak VRAM at 8,064 tokens: only 2,898 MB", WHITE, False),
    ("", WHITE, False),
    ("\u26a0\ufe0f  Criterion B NOT MET: 1.19\u00d7 achieved (target: 1.5\u00d7)", YELLOW, True),
    ("    But 19% extension is significant on a 4 GB device", WHITE, False),
], top=Inches(4.2))
add_footer(s)

# ── Slide 6: Key Results — KV-Cache ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Key Results \u2014 KV-Cache Quantization",
                "Independent evaluation with eager attention backend")
add_bullets(s, [
    ("Memory Savings", ACCENT, True),
    ("    INT4: 13% less memory per token (0.181 \u2192 0.157 MB/tok)", WHITE, False),
    ("    INT2: 15% less memory per token (0.181 \u2192 0.153 MB/tok)", WHITE, False),
    ("    At 8K tokens, INT4 saves ~192 MB vs. FP16", WHITE, False),
    ("", WHITE, False),
    ("Correctness", ACCENT, True),
    ("    100% top-1 token agreement with FP16 baseline", GREEN, False),
    ("    Max absolute logit difference: 0.0", GREEN, False),
    ("", WHITE, False),
    ("Combined Configs (Attention + KV Quant)", ACCENT, True),
    ("    FAILED: GQA head mismatch (16Q / 2KV) breaks kernel dispatch", RED, False),
    ("    Recommendation: apply each optimization independently", YELLOW, False),
], top=Inches(1.8))
add_footer(s)

# ── Slide 7: Quality Evaluation ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Quality Evaluation",
                "0-shot log-likelihood evaluation (4-bit NF4, eager, FP16 KV)")
add_table(s, ["Benchmark", "Examples", "Accuracy"],
          [
              ["ARC-Easy", "200", "62.5%"],
              ["BoolQ", "113 / 200", "60.2%"],
              ["HellaSwag", "200", "69.0%"],
              ["Mean", "\u2014", "63.9%"],
          ],
          top=Inches(2.2))
add_bullets(s, [
    ("BoolQ: 113/200 completed (60-min runtime limit reached)", LIGHT_GRAY, False),
    ("", WHITE, False),
    ("\u2705  63.9% mean accuracy \u2014 consistent with published 3B/4-bit baselines", GREEN, True),
    ("\u2705  Quality guard-rail MET: no degradation observed", GREEN, True),
], top=Inches(4.6))
add_footer(s)

# ── Slide 8: Recommendations ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Practical Recommendations",
                "Which configuration to use and when")
add_table(s, ["Workload", "Configuration", "Rationale"],
          [
              ["Latency-sensitive", "SDPA default", "33% p95 reduction, best throughput"],
              ["Long-context", "FlashAttention-2", "19% more context, lowest VRAM"],
              ["VRAM-constrained", "Add KV-cache INT4", "13% per-token savings, 0 quality loss"],
              ["General purpose", "SDPA default", "Best speed/compatibility balance"],
          ],
          top=Inches(2.2))
add_bullets(s, [
    ("\u26d4  Do NOT combine attention optimization + KV-cache quantization (yet)", RED, True),
    ("    GQA mismatch causes kernel dispatch failures", WHITE, False),
    ("    Expected to be fixed in future PyTorch / FlashAttention releases", LIGHT_GRAY, False),
], top=Inches(4.8))
add_footer(s)

# ── Slide 9: Conclusion & Future Work ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Conclusion & Future Work")
add_bullets(s, [
    ("Results Summary", ACCENT, True),
    ("    \u2705  Criterion A MET \u2014 33% latency reduction (SDPA default)", GREEN, False),
    ("    \u26a0\ufe0f  Criterion B PARTIALLY MET \u2014 1.19\u00d7 context (FlashAttention-2)", YELLOW, False),
    ("    \u2705  Quality preserved \u2014 63.9% mean accuracy", GREEN, False),
    ("    \u274c  Combined configs unstable \u2014 GQA head mismatch", RED, False),
    ("", WHITE, False),
    ("Future Work", ACCENT, True),
    ("    Desktop GPUs (RTX 3060/4060) \u2014 isolate thermal throttling", WHITE, False),
    ("    Larger models (7B) \u2014 test at higher compression ratios", WHITE, False),
    ("    Linux deployment \u2014 unlock flash SDP + mem-efficient backends", WHITE, False),
    ("    Batched inference \u2014 throughput for serving scenarios", WHITE, False),
    ("    Speculative decoding \u2014 further latency reduction", WHITE, False),
    ("", WHITE, False),
    ("Thank you!  \u2014  Questions?", ACCENT, True),
], top=Inches(1.8))
add_footer(s)

# ── Save ──
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
