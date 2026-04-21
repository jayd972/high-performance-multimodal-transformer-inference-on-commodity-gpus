"""
Generate a 3-slide weekly progress PPT.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Weekly_Progress_Weeks1-3.pptx")

# ── Colors ──
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCA, 0xCA, 0xDA)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)
RED_LIGHT = RGBColor(0xE7, 0x4C, 0x3C)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.4), left=Inches(0.6)):
    txBox = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    return txBox


def add_subtitle(slide, text, top=Inches(1.1), left=Inches(0.6)):
    txBox = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT
    return txBox


def add_bullet_block(slide, items, top=Inches(1.7), left=Inches(0.7), width=Inches(8.4)):
    """items: list of (text, color_or_None) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(2)

        run = p.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = color if color else WHITE

    return txBox


def add_table(slide, headers, rows, top=Inches(3.4), left=Inches(0.7)):
    """Add a small styled table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top,
                                        Inches(8.4), Inches(0.35 * num_rows))
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x24, 0x24, 0x3E)

    return tbl_shape


def add_footer(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

FOOTER = "Efficient Transformer Inference on Commodity GPUs  |  Jay Darji & Karma Patel"

# ──────────────────────────────────────────────────────────────
# SLIDE 1 — Achievements
# ──────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
set_slide_bg(s1, DARK_BG)

add_title(s1, "Achievements  (Weeks 1 - 3)")
add_subtitle(s1, "Planning, Environment Setup & Baseline Measurement")

achievements = [
    ("Week 1 - Planning & Scope Definition", ACCENT),
    ("    Selected primary model: Qwen2.5-3B-Instruct (3B params, ~1.87 GB at 4-bit)", WHITE),
    ("    Defined VRAM budget: 3.2 GB weights / 4.0 GB total on RTX 3050 Ti", WHITE),
    ("    Fixed all benchmark settings: prompt lengths [128, 512, 1024], greedy decoding", WHITE),
    ("    Documented correctness & quality evaluation protocols", WHITE),
    ("", None),
    ("Week 2 - Environment Setup & Validation", ACCENT),
    ("    Validated: Python 3.12, PyTorch 2.6.0+cu124, CUDA 12.4, bitsandbytes 0.49", WHITE),
    ("    Model loads in ~60s, weight footprint 1.87 GB (within 3.2 GB budget)", WHITE),
    ("    GPU/CPU memory monitoring harness tested and operational", WHITE),
    ("", None),
    ("Week 3 - Baseline Measurement", ACCENT),
    ("    Completed benchmark sweep across 3 prompt lengths (batch=1)", WHITE),
    ("    OOM threshold search: stable up to 4096 tokens (spills to shared memory beyond 4 GB)", WHITE),
    ("    Collected 5 deterministic correctness traces (20 steps each) as baseline reference", WHITE),
]
add_bullet_block(s1, achievements, top=Inches(1.6))

# Benchmark table
headers = ["Prompt Len", "p50 Latency", "p95 Latency", "Tokens/s", "Peak VRAM"]
rows = [
    ["128 tok", "29,764 ms", "31,449 ms", "2.13", "2,026 MB"],
    ["512 tok", "34,614 ms", "38,486 ms", "1.86", "2,061 MB"],
    ["1024 tok", "40,242 ms", "47,813 ms", "1.60", "2,195 MB"],
]
add_table(s1, headers, rows, top=Inches(5.8), left=Inches(0.7))

add_footer(s1, FOOTER)


# ──────────────────────────────────────────────────────────────
# SLIDE 2 — Planned Tasks for Next Week
# ──────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
set_slide_bg(s2, DARK_BG)

add_title(s2, "Planned Tasks  (Week 4)")
add_subtitle(s2, "Profiling & Attention Backend Planning")

plans = [
    ("1.  Profile baseline inference with PyTorch Profiler", GREEN),
    ("      Capture CUDA kernel timings, identify top time & memory contributors", WHITE),
    ("      Export Chrome-compatible trace for manual inspection", WHITE),
    ("", None),
    ("2.  Categorize bottleneck kernels", GREEN),
    ("      Group into: attention, linear/matmul, activation, normalization, memory ops", WHITE),
    ("      Compute percentage of total GPU time per category", WHITE),
    ("", None),
    ("3.  Plan attention backend strategy", GREEN),
    ("      Test SDPA backend availability: flash, memory-efficient, math", WHITE),
    ("      Determine primary backend and fallback path for RTX 3050 Ti", WHITE),
    ("      Add instrumentation to log which backend actually runs", WHITE),
    ("", None),
    ("4.  Deliverables", ACCENT),
    ("      profiling_results.json  -  Top 20 CUDA kernels and CPU operators", WHITE),
    ("      bottleneck_analysis.json  -  Categorized analysis with optimization recommendations", WHITE),
    ("      attention_backend_plan.json  -  Primary/fallback backend plan with evidence", WHITE),
    ("      Chrome profiler trace  -  For manual deep-dive inspection", WHITE),
]
add_bullet_block(s2, plans, top=Inches(1.6))

add_footer(s2, FOOTER)


# ──────────────────────────────────────────────────────────────
# SLIDE 3 — Issues & Comments
# ──────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
set_slide_bg(s3, DARK_BG)

add_title(s3, "Issues & Comments")
add_subtitle(s3, "Observations, risks, and mitigations")

issues = [
    ("Observation: Shared GPU Memory Spill", YELLOW),
    ("    RTX 3050 Ti does not hard-OOM at 4 GB; instead it spills to system RAM via", WHITE),
    ("    CUDA unified memory. This avoids crashes but causes severe slowdown:", WHITE),
    ("    throughput drops from ~6 tok/s (3072 tokens) to ~1.15 tok/s (4096 tokens).", WHITE),
    ("    We will report the \"practical OOM\" as the point where throughput drops >50%.", WHITE),
    ("", None),
    ("Observation: Slow Inference on Laptop GPU", YELLOW),
    ("    Baseline generates at ~2 tok/s for 64 output tokens. Full benchmark sweeps", WHITE),
    ("    take 15-30 minutes per configuration. We reduced benchmark runs from 10 to 5", WHITE),
    ("    and output limit from 128 to 64 tokens to keep iteration time manageable.", WHITE),
    ("", None),
    ("Risk: FlashAttention Compatibility", RED_LIGHT),
    ("    FlashAttention-2 (external package) may not install on Windows + RTX 3050 Ti.", WHITE),
    ("    Mitigation: PyTorch SDPA has built-in flash and memory-efficient backends.", WHITE),
    ("    We will test all three SDPA backends in Week 4 and fall back as needed.", WHITE),
    ("", None),
    ("Status: On Track", GREEN),
    ("    All Week 1-3 deliverables complete. Codebase is modular and reproducible.", WHITE),
    ("    Each weekly script runs independently and saves results to versioned folders.", WHITE),
]
add_bullet_block(s3, issues, top=Inches(1.6))

add_footer(s3, FOOTER)

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
