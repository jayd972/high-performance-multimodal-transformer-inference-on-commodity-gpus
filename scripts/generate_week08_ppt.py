"""
Generate Week 8 progress presentation — 3 slides only.
Slide 1-2: Week 8 results, Slide 3: Week 9 plan.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

import os, sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import RESULTS_DIR

OUTPUT = os.path.join(RESULTS_DIR, "Weekly_Progress_Week08.pptx")

DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
TITLE_BG   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
YELLOW     = RGBColor(0xF1, 0xC4, 0x0F)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ROW_DARK   = RGBColor(0x15, 0x25, 0x3D)
ROW_LIGHT  = RGBColor(0x1A, 0x2E, 0x4A)

FOOTER_TEXT = "LLM Inference Optimization  |  Darji & Patel  |  Week 8"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, items, top=Inches(1.8), left=Inches(0.7), width=Inches(8.6)):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, clr, bld) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.color.rgb = clr
        run.font.bold = bld
    return box


def add_table(slide, headers, rows, top=Inches(2.8), left=Inches(0.5),
              width=Inches(9.0)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.35)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   row_h * n_rows)
    tbl = shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    for i, row in enumerate(rows):
        bg = ROW_DARK if i % 2 == 0 else ROW_LIGHT
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
    return shape


def add_footer(slide):
    add_textbox(slide, FOOTER_TEXT,
                Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3),
                font_size=8, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, title,
                Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.7),
                font_size=26, bold=True, color=ACCENT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.4),
                    font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Slide 1: Week 8 Results — Tradeoff Analysis ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 8 Results — KV-Cache Tradeoff Analysis",
                "FP16 vs INT4 vs INT2 across prompt lengths 128 / 512 / 1024")

add_table(s, ["Prompt", "KV Type", "FP16 VRAM", "Quant VRAM", "Saved", "FP16 Latency", "Quant Latency", "Lat. Change"],
          [
              ["128",  "INT4", "2,026 MB", "2,023 MB", "0.2%", "12,785 ms", "15,282 ms", "+19.5%"],
              ["128",  "INT2", "2,026 MB", "2,022 MB", "0.2%", "12,785 ms", "16,060 ms", "+25.6%"],
              ["512",  "INT4", "2,061 MB", "2,048 MB", "0.6%", "16,767 ms", "16,503 ms", "\u22121.6%"],
              ["512",  "INT2", "2,061 MB", "2,046 MB", "0.7%", "16,767 ms", "17,118 ms", "+2.1%"],
              ["1024", "INT4", "2,195 MB", "2,171 MB", "1.1%", "21,326 ms", "21,899 ms", "+2.7%"],
              ["1024", "INT2", "2,195 MB", "2,167 MB", "1.3%", "21,326 ms", "22,233 ms", "+4.3%"],
          ],
          top=Inches(1.6), left=Inches(0.2), width=Inches(9.6))

add_bullets(s, [
    ("Key Observations", ACCENT, True),
    ("  VRAM savings grow with context: 0.2% at 128 tok \u2192 1.1\u20131.3% at 1024 tok", WHITE, False),
    ("  INT4 latency overhead is modest at longer prompts (+2.7% at 1024)", GREEN, False),
    ("  INT2 saves slightly more VRAM but with higher latency cost (+4.3%)", YELLOW, False),
    ("  At 128 tokens, both show warm-up artifacts (high latency overhead)", LIGHT_GRAY, False),
], top=Inches(5.0))

add_footer(s)

# ── Slide 2: Week 8 Results — Quality Eval & Final Selection ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 8 Results — Quality Evaluation & Final KV Selection",
                "0-shot log-likelihood accuracy (4-bit NF4 model, eager attention, FP16 KV)")

add_table(s, ["Benchmark", "Examples", "Correct", "Accuracy", "Runtime"],
          [
              ["ARC-Easy", "200", "125", "62.5%", "10.5 min"],
              ["BoolQ", "113 / 200", "68", "60.2%", "121 min"],
              ["HellaSwag", "200", "138", "69.0%", "16.2 min"],
              ["Mean", "\u2014", "\u2014", "63.9%", "\u2014"],
          ],
          top=Inches(1.7), left=Inches(1.0), width=Inches(8.0))

add_bullets(s, [
    ("BoolQ: only 113/200 completed (60-min limit reached per config)", LIGHT_GRAY, False),
    ("63.9% mean — consistent with published 3B / 4-bit baselines", GREEN, True),
], top=Inches(3.6))

add_bullets(s, [
    ("Final KV-Cache Selection", ACCENT, True),
], top=Inches(4.5))

add_table(s, ["Selected", "Avg VRAM Saved", "Avg Latency Change", "Rationale"],
          [
              ["INT4", "~0.6%", "+6.9%", "Best balance of savings and performance"],
          ],
          top=Inches(5.1), left=Inches(0.8), width=Inches(8.4))

add_bullets(s, [
    ("Alternative: INT2 saves ~0.7% VRAM but +10.7% latency — too costly", YELLOW, False),
    ("INT4 recommended for Week 9 combined configuration testing", GREEN, False),
], top=Inches(5.9))

add_footer(s)

# ── Slide 3: Week 9 Plan ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_BG)
add_slide_title(s, "Week 9 Plan — Combined Configuration Testing")

add_bullets(s, [
    ("Task 1 — Stability Test", ACCENT, True),
    ("  Combine each attention backend with KV-cache INT4 quantization", WHITE, False),
    ("  Test: SDPA default + INT4, FlashAttention-2 + INT4, mem-efficient + INT4", WHITE, False),
    ("  Record pass/fail status and any error messages", WHITE, False),
    ("", WHITE, False),
    ("Task 2 — Combined Benchmarks", ACCENT, True),
    ("  For stable combinations: run full benchmark (128/512/1024 tokens)", WHITE, False),
    ("  Measure latency, throughput, and peak VRAM", WHITE, False),
    ("  Compare against individual optimization results", WHITE, False),
    ("", WHITE, False),
    ("Task 3 — Incompatibility Report", ACCENT, True),
    ("  Document any failed combinations with root cause analysis", WHITE, False),
    ("  Especially GQA head mismatch issues (16Q vs 2KV)", WHITE, False),
    ("", WHITE, False),
    ("Task 4 — Final Benchmark Runbook", ACCENT, True),
    ("  Define the exact configuration matrix for Week 10 full benchmark", WHITE, False),
    ("  Only include configurations proven stable in this week", WHITE, False),
    ("", WHITE, False),
    ("Expected Deliverables", ACCENT, True),
    ("  combined_stability_test.json  |  combined_benchmark_results.json", LIGHT_GRAY, False),
    ("  incompatibility_report.json   |  final_benchmark_runbook.json", LIGHT_GRAY, False),
], top=Inches(1.5))

add_footer(s)

# ── Save ──
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
