"""Generate Week 06 results + Week 07 plan PowerPoint (3 slides)."""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config import RESULTS_DIR

PPT_PATH = os.path.join(RESULTS_DIR, "Weekly_Progress_Week6.pptx")

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT = RGBColor(0x89, 0xB4, 0xFA)
GREEN = RGBColor(0xA6, 0xE3, 0xA1)
YELLOW = RGBColor(0xF9, 0xE2, 0xAF)
WHITE = RGBColor(0xCD, 0xD6, 0xF4)
SUBTEXT = RGBColor(0xBA, 0xC2, 0xDE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, size=18,
              color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def _add_table(slide, left, top, width, rows_data):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                        Inches(0.38 * n_rows))
    tbl = tbl_shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    paragraph.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x3B, 0x3D, 0x54)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    RGBColor(0x2A, 0x2B, 0x3D) if r % 2 == 1
                    else RGBColor(0x33, 0x34, 0x48)
                )
    return tbl_shape


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ================================================================
    # SLIDE 1 - Week 6 Results: Benchmark Table
    # ================================================================
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1, DARK_BG)

    _add_text(s1, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 06 - Attention Optimization Final Results",
              size=28, color=ACCENT, bold=True)
    _add_text(s1, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
              "Model: Qwen/Qwen2.5-3B-Instruct  |  GPU: RTX 3050 Ti (4 GB)  |  batch=1, 64 new tokens",
              size=13, color=SUBTEXT)

    _add_text(s1, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
              "Benchmark Results (prompt lengths 128, 512, 1024)",
              size=16, color=YELLOW, bold=True)

    bench_rows = [
        ["Backend", "128 tok", "512 tok", "1024 tok", "VRAM (MB)", "OOM Threshold"],
        ["Eager",           "12.7s / 5.0",  "24.4s / 2.6",  "29.9s / 2.1",  "2,027-2,153", "7,168"],
        ["SDPA Default",    "36.5s / 1.8",  "41.5s / 1.5",  "59.1s / 1.0",  "3,397-3,565", "4,096"],
        ["SDPA Math",      "50.5s / 1.3",  "43.8s / 1.5",  "56.9s / 1.1",  "3,397-3,565", "-"],
        ["FlashAttention-2","39.9s / 1.6", "40.9s / 1.6",  "14.3s / 4.5*", "3,398-3,480", "4,096"],
    ]
    _add_table(s1, Inches(0.6), Inches(1.9), Inches(11.5), bench_rows)

    _add_text(s1, Inches(0.6), Inches(4.2), Inches(11), Inches(0.5),
              "* FA2 at 1024: anomalously fast (thermal recovery after model reload). "
              "Format: p50 latency / tok/s",
              size=10, color=SUBTEXT)

    _add_text(s1, Inches(0.6), Inches(4.9), Inches(5), Inches(0.4),
              "Key Findings",
              size=16, color=YELLOW, bold=True)
    _add_text(s1, Inches(0.6), Inches(5.3), Inches(11), Inches(1.2),
              "- Eager uses ~1.4x less VRAM (2.0 vs 3.4 GB) and supports longer context (7k vs 4k tokens)\n"
              "- Thermal throttling: later runs 2-3x slower; allow warmup before benchmarking\n"
              "- Correctness vs eager baseline: 100% token agreement for all backends",
              size=12, color=GREEN)

    # ================================================================
    # SLIDE 2 - Week 6 Summary & Deliverables
    # ================================================================
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2, DARK_BG)

    _add_text(s2, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 06 - Deliverables & Stability Notes",
              size=28, color=ACCENT, bold=True)

    _add_text(s2, Inches(0.6), Inches(1.1), Inches(5), Inches(0.4),
              "Deliverables",
              size=16, color=YELLOW, bold=True)
    _add_text(s2, Inches(0.6), Inches(1.5), Inches(11), Inches(1.0),
              "attention_final_results.json/csv  |  attention_oom_thresholds.json  |  "
              "stability_notes.json  |  report_attention_section.md  |  discrepancy_audit.md",
              size=12, color=SUBTEXT)

    _add_text(s2, Inches(0.6), Inches(2.7), Inches(5), Inches(0.4),
              "Plots (results/plots/week06_attention/)",
              size=16, color=YELLOW, bold=True)
    _add_text(s2, Inches(0.6), Inches(3.1), Inches(11), Inches(0.5),
              "attention_latency_comparison.png  |  attention_throughput_comparison.png  |  attention_vram_comparison.png",
              size=12, color=SUBTEXT)

    _add_text(s2, Inches(0.6), Inches(3.9), Inches(5), Inches(0.4),
              "Backend Availability (Windows)",
              size=16, color=YELLOW, bold=True)
    avail_rows = [
        ["Backend", "Status", "Notes"],
        ["SDPA Flash",         "Unavailable", "Not in PyTorch Windows wheel"],
        ["SDPA Mem-Efficient", "Fails (GQA)", "16Q vs 2KV head mismatch on Qwen2.5-3B"],
        ["SDPA Math",          "OK",         "Universal fallback"],
        ["FlashAttention-2",   "OK",         "flash-attn pre-built wheel"],
    ]
    _add_table(s2, Inches(0.6), Inches(4.3), Inches(10), avail_rows)

    # ================================================================
    # SLIDE 3 - Week 7 Plan: KV-Cache Quantization
    # ================================================================
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3, DARK_BG)

    _add_text(s3, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
              "Week 07 Plan - KV-Cache Quantization Phase",
              size=28, color=ACCENT, bold=True)

    _add_text(s3, Inches(0.6), Inches(1.0), Inches(11), Inches(0.4),
              "Objective: Prototype KV-cache quantization (INT8/INT4), measure memory savings and correctness, decide on deeper integration.",
              size=15, color=SUBTEXT)

    tasks = [
        ("1. KV-Cache Quantization Support Check",
         "Verify transformers + quanto support for KV-cache quantization. Document "
         "available methods (fp16, int8, int4) and API usage."),

        ("2. End-to-End Prototype",
         "Run inference with INT8 and INT4 KV-cache quantization. Measure latency, "
         "peak VRAM, and verify model generates valid output."),

        ("3. Correctness Verification",
         "Compare token outputs and logits against Week 03 baseline. Report agreement "
         "rate and max logit difference per quantization type."),

        ("4. Memory per Token Measurement",
         "Measure VRAM growth vs prompt length for each KV type (fp16, int8, int4). "
         "Quantify memory savings from quantization."),

        ("5. Integration Decision",
         "Document whether to pursue deeper framework integration for Week 08 based "
         "on prototype results, correctness, and memory savings."),
    ]

    y = Inches(1.6)
    for title, desc in tasks:
        _add_text(s3, Inches(0.8), y, Inches(10), Inches(0.35),
                  title, size=16, color=GREEN, bold=True)
        _add_text(s3, Inches(1.0), y + Inches(0.36), Inches(10), Inches(0.5),
                  desc, size=12, color=SUBTEXT)
        y += Inches(0.95)

    _add_text(s3, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
              "Deliverables: kv_quant_support.json  |  kv_quant_prototype.json  |  kv_correctness_report.json  |  kv_memory_per_token.json  |  integration_decision.json",
              size=12, color=YELLOW, bold=True)

    prs.save(PPT_PATH)
    print(f"Saved: {PPT_PATH}")


if __name__ == "__main__":
    build_ppt()
